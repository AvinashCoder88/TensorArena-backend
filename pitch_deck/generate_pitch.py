from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Define content
    slides_content = [
        {
            "layout": 0, # Title Slide
            "title": "TensorArena",
            "subtitle": "The Future of AI-Adaptive Technical Interview Preparation\n\nPresenter: [Your Name]"
        },
        {
            "layout": 1, # Title and Content
            "title": "The Problem: The Broken Interview Prep Loop",
            "content": [
                "Disconnect: Standard coding platforms focus on syntax, not engineering reality.",
                "Generic: One-size-fits-all questions don't screen for specialized roles like MLE or System Design.",
                "Static: No feedback on code style, efficiency, or architectural choices—only 'Pass/Fail'.",
                "Expensive: Mock interviews cost hundreds of dollars per session."
            ]
        },
        {
            "layout": 1, 
            "title": "The Solution: TensorArena",
            "content": [
                "Role-Specific: Dedicated tracks for ML Engineers, Data Scientists, and AI Product Managers.",
                "Beyond Algorithms: Modules for System Design, Paper Implementation, and Production Incident Simulation.",
                "AI Tutor: Real-time feedback on Correctness, Efficiency (Big O), and Code Style.",
                "Voice-First: Mock Interview mode with interactive AI personas."
            ]
        },
        {
            "layout": 1,
            "title": "Key Features: Comprehensive Training Ground", 
            "content": [
                "The Arena: Adaptive Coding Challenges (Python, Algo, MLE) with specialized modes.",
                "System Design Mode: Architecture-focused challenges.",
                "Role-Play Arena: Real-world scenarios (e.g., 'Debug this Prod Incident').",
                "Mock Interview: Voice-interactive technical interviews with AI personas.",
                "AI Feedback Loop: Instant detailed grading on Code Quality & Complexity."
            ]
        },
        {
            "layout": 1,
            "title": "Market Opportunity",
            "content": [
                "Target Audience: 25M+ Software Developers worldwide.",
                "High-Growth Niche: AI/ML specialization segment is expanding rapidly.",
                "B2C: Job seekers optimizing for FAANG+ roles.",
                "B2B: Companies using TensorArena for initial candidate screening.",
                "B2E (Education): Bootcamps and Universities using 'Classroom Mode'."
            ]
        },
        {
            "layout": 1,
            "title": "Student-Teacher Automation (New!)",
            "content": [
                "Automated Grading: AI instantaneously grades code for correctness, style, and efficiency—saving teachers hundreds of hours.",
                "Classroom Dashboard: Teachers can track progress, identify struggling students, and manage assignments in real-time.",
                "Curriculum Integration: Custom problem sets aligned with course syllabus.",
                "Scalable Feedback: Every student gets personalized 1:1 mentorship from the AI."
            ]
        },
        {
            "layout": 1,
            "title": "Business Model: Freemium SaaS",
            "content": [
                "Free Tier: 5 complimentary usage credits (Try before you buy).",
                "Pro Subscription ($19.99/mo):",
                "  - Unlimited Questions",
                "  - System Design & Role-Based Modes",
                "  - Deep AI Analytics & Progress Tracking", 
                "Enterprise: Custom role definitions and candidate analytics dashboard (Planned)."
            ]
        },
        {
            "layout": 1,
            "title": "Technology Stack: Built for Scale",
            "content": [
                "Frontend: Next.js 14, React, Tailwind CSS, Monaco Editor.",
                "Backend: FastAPI (Python), Prisma ORM for robust data management.",
                "AI Core: Integration with OpenAI GPT-4 & Google Gemini for adaptive intelligence.",
                "Real-Time: Deepgram (Voice/Audio processing) & WebSockets for live interaction."
            ]
        },
        {
            "layout": 1,
            "title": "Roadmap: Vision for the Future",
            "content": [
                "Q3 2024: Multiplayer Battle Mode (1v1 Coding Duels).",
                "Q4 2024: Corporate Integration (ATS connect).",
                "Q1 2025: Full Voice-Native System Design Interviews.",
                "Q2 2025: Expanded Role Library (DevOps, SRE, Frontend)."
            ]
        },
        {
            "layout": 0, # Title Slide (Conclusion)
            "title": "Join the Revolution",
            "subtitle": "TensorArena is redefining how engineers prepare and hire.\n\nContact: [Your Contact Info]"
        }
    ]

    for slide_data in slides_content:
        # 0 is usually Title, 1 is Title and Content
        layout_index = slide_data.get("layout", 1) 
        
        # Adjust layout based on keys if needed, but 1 is safe for lists
        try:
            slide_layout = prs.slide_layouts[layout_index]
        except:
            slide_layout = prs.slide_layouts[1] # Fallback
            
        slide = prs.slides.add_slide(slide_layout)

        # Set Title
        if "title" in slide_data:
            try:
                title = slide.shapes.title
                title.text = slide_data["title"]
            except:
                pass

        # Set Content/Subtitle
        if layout_index == 0: # Title Slide
            if "subtitle" in slide_data:
                try:
                    subtitle = slide.placeholders[1]
                    subtitle.text = slide_data["subtitle"]
                except:
                    pass
        else: # Content Slide
            if "content" in slide_data:
                try:
                    # Depending on template, placeholder indices vary. 
                    # Usually title is 0, body is 1.
                    body_shape = slide.shapes.placeholders[1]
                    tf = body_shape.text_frame
                    tf.clear() # Clear default prompt
                    
                    for i, item in enumerate(slide_data["content"]):
                        p = tf.add_paragraph()
                        p.text = item
                        p.level = 0
                        
                        if item.strip().startswith("-"):
                             p.text = item.replace("-", "").strip()
                             p.level = 1
                except:
                    pass

    # Save
    output_file = "TensorArena_PitchDeck.pptx"
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    create_presentation()
